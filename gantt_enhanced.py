from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def rect(slide, x, y, w, h, fill, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    if line: s.line.color.rgb = rgb(line); s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, sz=10, bold=False, col="FFFFFF",
        align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.name = font; r.font.color.rgb = rgb(col)
    return tb

# Colors
BG="0B0F1A"; ACC="3B82F6"; WHT="FFFFFF"; SLV="94A3B8"; BDR="1E293B"
TRN="4F46E5"; CUR="2563EB"; PLN="7C3AED"; RES="BE185D"; MIL="F59E0B"
DK1="0D1220"; DK2="111827"; DK3="0F172A"
PHASE1="1E40AF"; PHASE2="059669"; PHASE3="DC2626"; PHASE4="7C3AED"

prs = Presentation()
prs.slide_width = Inches(13.3)
prs.slide_height = Inches(7.5)
BLK = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLK)
rect(s1, 0,0,13.3,7.5, BG)
rect(s1, 0,0,13.3,0.07, ACC)
rect(s1, 0,0.07,0.35,7.43, "0D1220")

txt(s1,"MY GRADUATE PROGRAM TIMELINE",0.6,1.3,10.5,1.0,sz=34,bold=True,col=WHT,font="Trebuchet MS")
txt(s1,"PhD in Technology  ·  Purdue Polytechnic Institute  ·  Spring 2026 – Fall 2029",0.6,2.48,11,0.4,sz=14,col=SLV)
txt(s1,"Md Kamruzzaman Kamrul",0.6,3.02,7,0.38,sz=14,bold=True,col=ACC)
txt(s1,"Bowen School of Construction Management Technology",0.6,3.45,9,0.32,sz=12,col=SLV)
txt(s1,"Advisor: Dr. Heung Jin Oh  ·  Wii Lab  ·  CSE Concentration",0.6,3.82,9,0.32,sz=12,col=SLV)

for i,(v,l) in enumerate([("103","TOTAL CREDITS"),("21","TRANSFER (GSU)"),("41","PURDUE COURSEWORK"),("41","DISSERTATION"),("4 YRS","DURATION")]):
    bx=0.6+i*2.42
    rect(s1,bx,4.55,2.2,1.2,DK2,BDR,0.5)
    txt(s1,v,bx,4.62,2.2,0.6,sz=30,bold=True,col=ACC,align=PP_ALIGN.CENTER,font="Trebuchet MS")
    txt(s1,l,bx,5.24,2.2,0.32,sz=9,col=SLV,align=PP_ALIGN.CENTER)

for i,(c,l) in enumerate([(TRN,"Transfer (GSU)"),(CUR,"Current (Sp 2026)"),(PLN,"Planned Courses"),(RES,"Dissertation Research"),(MIL,"Milestone")]):
    lx=0.6+i*2.42
    rect(s1,lx,6.38,0.22,0.22,c)
    txt(s1,l,lx+0.28,6.34,2.0,0.3,sz=10,col=SLV)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — GANTT CHART
# ══════════════════════════════════════════════════════════
# Semesters (12 semesters: Sp 26 → Fa 29)
SEMS=["Sp'26","Su'26","Fa'26","Sp'27","Su'27","Fa'27","Sp'28","Su'28","Fa'28","Sp'29","Su'29","Fa'29"]
N=len(SEMS)

# Rows: (short_label, full_label, group, start, span, color)
ROWS=[
 # Transfer
 ("GSU Transfer","GSU MS Transfer — 7 Courses (21 cr)","TRANSFER",0,1,TRN),
 # Core Required
 ("TECH 60300","TECH 60300  Grad Seminar Planning (1cr)","CORE REQUIRED",0,1,CUR),
 ("MET 52700","MET 52700  Technology: Global Perspective (3cr)","CORE REQUIRED",2,1,PLN),
 ("STAT 50100","STAT 50100  Experimental Statistics I (3cr)","CORE REQUIRED",2,1,PLN),
 ("TECH 60100","TECH 60100  Research Seminar (1cr)","CORE REQUIRED",2,1,PLN),
 ("TECH 64600","TECH 64600  Analysis of Research in Industry (3cr)","CORE REQUIRED",2,1,PLN),
 ("TECH 67600","TECH 67600  Analysis of Research Methods (3cr)","CORE REQUIRED",3,1,PLN),
 # Technology/Discovery
 ("STAT 52000","STAT 52000  Time Series & Applications (3cr)","TECH / DISCOVERY",0,1,CUR),
 ("STAT 52400","STAT 52400  Applied Multivariate Analysis (3cr)","TECH / DISCOVERY",5,1,PLN),
 ("STAT 51400","STAT 51400  Design of Experiments (3cr)","TECH / DISCOVERY",5,1,PLN),
 # CSE Concentration
 ("CS 57800","CS 57800  Statistical Machine Learning (3cr)","CSE CONCENTRATION",3,1,PLN),
 ("CS 52000","CS 52000  Computational Optimization (3cr)","CSE CONCENTRATION",3,1,PLN),
 ("ECE 57000","ECE 57000  Artificial Intelligence (3cr)","CSE CONCENTRATION",5,1,PLN),
 ("GRAD 68900","GRAD 68900  CSE Seminar 0cr (recurring)","CSE CONCENTRATION",8,4,"1E3A5F"),
 # Cognate Electives
 ("CS 59300","CS 59300  Computer Vision (3cr)","COGNATE ELECTIVES",0,1,CUR),
 ("ECE 59500","ECE 59500  Selected Topics in ECE (3cr)","COGNATE ELECTIVES",0,1,CUR),
 ("CS 58000","CS 58000  Algorithm Design & Analysis (3cr)","COGNATE ELECTIVES",6,1,PLN),
 # Dissertation
 ("TECH 69900","TECH 69900  Dissertation Research: Phases 1-6","DISSERTATION",2,10,RES),
]

MILESTONES=[(5,"Prelim\nExam"),(8,"Proposal\nDefense"),(8,"Data\nCollection"),(10,"Writing"),(11,"Final\nDefense")]

s2 = prs.slides.add_slide(BLK)
rect(s2,0,0,13.3,7.5,BG)
rect(s2,0,0,13.3,0.07,ACC)
txt(s2,"MY GRADUATE PROGRAM TIMELINE",0.3,0.1,9,0.38,sz=17,bold=True,col=WHT,font="Trebuchet MS")
txt(s2,"Kamruzzaman Kamrul  ·  PhD Technology  ·  Purdue  ·  Spring 2026 – Fall 2029",0.3,0.5,11,0.24,sz=9,col=SLV)

# Layout
LBL_W = 2.9
CHART_X = LBL_W
CHART_W = 13.3 - LBL_W - 0.08
CELL_W = CHART_W / N
HDR_Y = 0.78
HDR_H = 0.28
ROW_Y0 = HDR_Y + HDR_H + 0.03
ROW_H = 0.195
ROW_GAP = 0.008
GRP_GAP = 0.055

# Compute y positions
row_ys = []
cy = ROW_Y0
prev_g = None
for (sl,fl,g,st,sp,col) in ROWS:
    if prev_g is not None and g != prev_g:
        cy += GRP_GAP
    row_ys.append(cy)
    cy += ROW_H + ROW_GAP
    prev_g = g
total_rows_h = cy - ROW_Y0

# Group map
grp_map = {}
for i,(sl,fl,g,st,sp,col) in enumerate(ROWS):
    grp_map.setdefault(g,[]).append(i)

# Semester column BGs + headers
for i,sem in enumerate(SEMS):
    cx = CHART_X + i*CELL_W
    is_su = sem.startswith("Su")
    is_cur = (i==0)
    cbg = "090D18" if is_su else ("0B0F1A" if i%2==0 else "0D1322")
    rect(s2, cx, ROW_Y0, CELL_W, total_rows_h+0.02, cbg)
    if is_cur:
        rect(s2, cx, ROW_Y0, CELL_W, total_rows_h+0.02, "0B1E38")
    hbg = ACC if is_cur else ("1C2D3D" if is_su else "1E293B")
    rect(s2, cx+0.01, HDR_Y+0.02, CELL_W-0.02, HDR_H-0.04, hbg)
    txt(s2, sem, cx, HDR_Y+0.03, CELL_W, HDR_H-0.05, sz=7.5,
        bold=is_cur, col=WHT if is_cur else SLV, align=PP_ALIGN.CENTER)

# Draw group labels + row labels + bars
drawn_grps = set()
for i,(sl,fl,g,st,sp,col) in enumerate(ROWS):
    ry = row_ys[i]

    # Group panel (once per group)
    if g not in drawn_grps:
        idxs = grp_map[g]
        g_top = row_ys[idxs[0]]
        g_bot = row_ys[idxs[-1]] + ROW_H
        g_h = g_bot - g_top
        strip = CUR if col==CUR else (RES if col==RES else (TRN if col==TRN else PLN))
        rect(s2, 0, g_top, 0.06, g_h, strip)
        rect(s2, 0.06, g_top, LBL_W-0.1, g_h, DK1)
        txt(s2, g, 0.1, g_top + g_h/2 - 0.11, LBL_W-0.18, 0.22,
            sz=6.5, bold=True, col="64748B")
        drawn_grps.add(g)

    # Row label
    rect(s2, 0.06, ry, LBL_W-0.1, ROW_H, DK1)
    txt(s2, fl, 0.12, ry+0.015, LBL_W-0.18, ROW_H-0.03,
        sz=6.2, col=SLV, wrap=False)

    # Bar
    bx = CHART_X + st*CELL_W + 0.015
    bw = sp*CELL_W - 0.03
    bar_h = ROW_H - 0.03
    rect(s2, bx, ry+0.015, bw, bar_h, col)

# Milestone markers
ms_base_y = row_ys[-1] + ROW_H + 0.08
for (sem_idx, ms_lbl) in MILESTONES:
    mx = CHART_X + sem_idx*CELL_W + CELL_W/2
    # Diamond
    d = s2.shapes.add_shape(1, Inches(mx-0.1), Inches(ms_base_y), Inches(0.2), Inches(0.2))
    d.fill.solid(); d.fill.fore_color.rgb = rgb(MIL)
    d.line.fill.background(); d.rotation = 45
    # Label
    txt(s2, ms_lbl, mx-0.45, ms_base_y+0.22, 0.9, 0.32,
        sz=6.5, bold=True, col=MIL, align=PP_ALIGN.CENTER, wrap=True)
    # Dotted line
    ln = s2.shapes.add_connector(1, Inches(mx), Inches(ROW_Y0), Inches(mx), Inches(ms_base_y+0.01))
    ln.line.color.rgb = rgb(MIL); ln.line.width = Pt(0.75); ln.line.dash_style = 4

# Legend
leg_y = 7.08
for i,(c,l) in enumerate([(TRN,"Transfer Credits"),(CUR,"Current (Sp'26)"),(PLN,"Planned Courses"),(RES,"Dissertation Research"),(MIL,"Milestone")]):
    lx = 0.3+i*2.55
    rect(s2,lx,leg_y+0.04,0.18,0.16,c)
    txt(s2,l,lx+0.24,leg_y,2.1,0.26,sz=9,col=SLV)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — MILESTONE DELIVERABLES (2-COLUMN GRID - FIXED)
# ══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLK)
rect(s3,0,0,13.3,7.5,BG)
rect(s3,0,0,13.3,0.07,ACC)
txt(s3,"4-YEAR PATH: PHASES & KEY DELIVERABLES",0.3,0.1,9,0.35,sz=18,bold=True,col=WHT,font="Trebuchet MS")
txt(s3,"Each milestone builds on prior coursework with specific outputs",0.3,0.48,11,0.22,sz=9,col=SLV)

# 2x2 Grid (2 columns, 2 rows)
milestones_grid = [
    {
        "phase": "PHASE 1: FOUNDATION",
        "dates": "Sp'26 – Sp'27",
        "color": PHASE1,
        "items": ["Front-load core & tech courses","Begin literature review (200+ papers)","Establish committee","Research problem identification"],
        "milestone": "✓ Research direction approved",
        "intensity": "8-15 hrs/week"
    },
    {
        "phase": "PHASE 2: QUALIFICATION",
        "dates": "Su'27 – Fa'27",
        "color": PHASE2,
        "items": ["Co-op internship (Jun-Aug)","Prelim exam preparation","Methodology development","Complete preliminary research outline"],
        "milestone": "★ PRELIM EXAM (Fa 27)",
        "intensity": "20+ hrs/week"
    },
    {
        "phase": "PHASE 3: PROPOSAL",
        "dates": "Sp'28 – Fa'28",
        "color": PHASE3,
        "items": ["Write proposal (30-50 pages)","Finalize research design & methods","Plan data collection strategy","Prepare proposal defense"],
        "milestone": "★ PROPOSAL DEFENSE (Fa 28)",
        "intensity": "18-20+ hrs/week"
    },
    {
        "phase": "PHASE 4: COMPLETION",
        "dates": "Sp'29 – Fa'29",
        "color": PHASE4,
        "items": ["Intensive data collection & analysis","Write Ch 1-4 (Intro, Lit, Methods, Results)","Dissertation composition","Final defense & graduation"],
        "milestone": "★ FINAL DEFENSE (Fa 29)",
        "intensity": "25+ hrs/week (peak)"
    },
]

col_w = 6.2
col_h = 3.2
start_x = [0.35, 6.8]
start_y = [0.8, 4.2]

box_idx = 0
for row in range(2):
    for col in range(2):
        if box_idx >= len(milestones_grid):
            break
        m = milestones_grid[box_idx]
        x = start_x[col]
        y = start_y[row]
        
        # Phase background
        rect(s3, x, y, col_w, 0.32, m["color"])
        txt(s3, m["phase"], x+0.15, y+0.04, col_w-0.3, 0.24, sz=9, bold=True, col=WHT)
        
        # Dates
        rect(s3, x, y+0.32, col_w, 0.22, DK2, BDR, 0.3)
        txt(s3, m["dates"], x+0.1, y+0.36, col_w-0.2, 0.18, sz=7, col=SLV)
        
        # Items
        item_y = y + 0.56
        for item in m["items"]:
            rect(s3, x, item_y, col_w, 0.16, DK3, BDR, 0.2)
            txt(s3, "• " + item, x+0.1, item_y+0.02, col_w-0.2, 0.14, sz=6.5, col=SLV)
            item_y += 0.17
        
        # Milestone (bottom)
        rect(s3, x, item_y+0.02, col_w, 0.22, MIL)
        txt(s3, m["milestone"], x+0.08, item_y+0.04, col_w-0.16, 0.18, sz=7, bold=True, col=WHT)
        
        box_idx += 1

# ══════════════════════════════════════════════════════════
# SLIDE 4 — MILESTONE REQUIREMENTS & RATIONALE
# ══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLK)
rect(s4,0,0,13.3,7.5,BG)
rect(s4,0,0,13.3,0.07,ACC)
txt(s4,"KEY MILESTONES: REQUIREMENTS & RATIONALE",0.3,0.1,9,0.35,sz=18,bold=True,col=WHT,font="Trebuchet MS")

milestones_req = [
    ("PRELIMINARY EXAM (Fall 2027)", MIL, [
        ("Requirement", "Written comprehensive exam + Oral defense + Committee approval"),
        ("Deliverables", "Pass both written + oral exams | Submit approved outline | Committee signs candidacy"),
        ("Timeline", "End of Year 2 (Semester 5) — Standard PhD progression"),
        ("Intensity", "20+ hours/week — Focused on exam prep & research design"),
    ]),
    ("PROPOSAL DEFENSE (Fall 2028)", PHASE3, [
        ("Requirement", "Defend 30-50 page proposal with full committee"),
        ("Deliverables", "Present & defend proposal | Submit approved document | Committee approves method"),
        ("Timeline", "12 months after Prelim — Allows methodology development"),
        ("Intensity", "18-20+ hours/week — Proposal writing + preliminary data collection"),
    ]),
    ("DATA COLLECTION & ANALYSIS (Sp-Su 2029)", PHASE4, [
        ("Requirement", "Acquire, validate, and analyze dissertation dataset(s)"),
        ("Deliverables", "Complete ≥1 validated dataset | Preprocessing pipeline documented | Results generated"),
        ("Timeline", "Semesters 9-10 — Peak research phase"),
        ("Intensity", "25+ hours/week (peak) — Full-time research focus"),
    ]),
    ("FINAL DEFENSE (Fall 2029)", MIL, [
        ("Requirement", "Oral defense & approval of completed dissertation"),
        ("Deliverables", "Final dissertation approved by committee | Pass oral defense | ETD submission"),
        ("Timeline", "End of Year 4 (Semester 11) — 4-year completion"),
        ("Intensity", "25+ hours/week — Final writing, revisions, defense prep"),
    ]),
]

y_pos = 0.6
for (title, col, reqs) in milestones_req:
    # Title bar
    rect(s4, 0.3, y_pos, 12.7, 0.28, col)
    txt(s4, title, 0.45, y_pos+0.04, 12.4, 0.22, sz=10, bold=True, col=WHT)
    y_pos += 0.3
    
    # Requirements table
    for (label, value) in reqs:
        rect(s4, 0.3, y_pos, 2.0, 0.18, DK3, BDR, 0.2)
        txt(s4, label, 0.4, y_pos+0.02, 1.8, 0.16, sz=7, bold=True, col=ACC)
        rect(s4, 2.35, y_pos, 10.65, 0.18, DK2, BDR, 0.2)
        txt(s4, value, 2.45, y_pos+0.02, 10.45, 0.16, sz=7, col=SLV, wrap=True)
        y_pos += 0.19
    
    y_pos += 0.05

# ══════════════════════════════════════════════════════════
# SLIDE 5 — COURSE INVENTORY
# ══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLK)
rect(s5,0,0,13.3,7.5,BG)
rect(s5,0,0,13.3,0.07,ACC)
txt(s5,"COURSE INVENTORY & CREDIT BREAKDOWN",0.3,0.12,11,0.44,sz=19,bold=True,col=WHT,font="Trebuchet MS")

sem_left=[
 ("TRANSFER — GSU MS (21 cr)",TRN,[
  "CIS 8005   Data Programming (3cr)",
  "CIS 8398   Advanced AI for Business (3cr)",
  "CIS 8695   Big Data Analytics (3cr)",
  "CIS 8795   Big Data Infrastructure (3cr)",
  "CIS 8080   IS Security & Privacy (3cr)",
  "CIS 8045   Unstructured Data Management (3cr)",
  "CIS 8690   Topics in Information Systems (3cr)",
 ]),
 ("SPRING 2026 — Current (10cr)",CUR,[
  "STAT 52000  Time Series & Applications (3cr)",
  "CS 59300    Computer Vision (3cr)",
  "ECE 59500   Selected Topics in ECE (3cr)",
  "TECH 60300  Grad Seminar – Planning (1cr)",
 ]),
 ("FALL 2026 (10cr)",PLN,[
  "MET 52700   Technology: Global Perspective (3cr)",
  "STAT 50100  Experimental Statistics I (3cr)",
  "TECH 60100  Research Seminar in Technology (1cr)",
  "TECH 64600  Analysis of Research in Industry (3cr)",
 ]),
 ("SPRING 2027 (9cr)",PLN,[
  "TECH 67600  Analysis of Research Methods (3cr)",
  "CS 57800    Statistical Machine Learning (3cr)",
  "CS 52000    Computational Optimization (3cr)",
 ]),
 ("SUMMER 2027 — CO-OP",RES,[
  "Industry Internship (3 months)",
 ]),
 ("FALL 2027 — ★ Prelim Exam",PLN,[
  "STAT 52400  Applied Multivariate Analysis (3cr)",
  "GRAD 68900  CSE Seminar (0cr)",
 ]),
]
sem_right=[
 ("SPRING 2028 — Proposal Writing (3cr)",PLN,[
  "STAT 51400  Design of Experiments (3cr)",
  "GRAD 68900  CSE Seminar (0cr)",
 ]),
 ("SUMMER 2028 — Proposal Dev",RES,[
  "TECH 69900  Dissertation (proposal draft)",
 ]),
 ("FALL 2028 — ★ Proposal Defense",RES,[
  "ECE 57000   Artificial Intelligence (3cr)",
  "GRAD 68900  CSE Seminar (0cr)",
 ]),
 ("SPRING 2029 — Data Collection",RES,[
  "CS 58000    Algorithm Design & Analysis (3cr)",
  "GRAD 68900  CSE Seminar (0cr)",
 ]),
 ("SUMMER 2029 — Analysis & Writing",RES,[
  "TECH 69900  Dissertation (analysis phase)",
 ]),
 ("FALL 2029 — ★ FINAL DEFENSE 🎓",MIL,[
  "TECH 69900  Dissertation Defense",
  "Graduation",
 ]),
]

def draw_col(slide, blocks, sx, sy, cw):
    y=sy; HH=0.235; RH=0.205
    for (title,col,courses) in blocks:
        rect(slide,sx,y,cw,HH,col)
        txt(slide,title,sx+0.07,y+0.025,cw-0.1,HH-0.04,sz=8.8,bold=True,col=WHT)
        y+=HH
        for c in courses:
            rect(slide,sx,y,cw,RH,DK2,BDR,0.3)
            txt(slide,c,sx+0.1,y+0.015,cw-0.15,RH-0.025,sz=8.2,col=SLV)
            y+=RH
        y+=0.05
    return y

draw_col(s5,sem_left,0.25,0.68,6.1)
draw_col(s5,sem_right,6.75,0.68,6.1)
rect(s5,0.25,6.82,12.8,0.52,DK3,BDR,0.5)
txt(s5,"CREDIT SUMMARY:  Transfer GSU 21cr  +  Purdue Coursework 41cr  +  Dissertation 41cr  =  103cr Listed  |  Graduation: Fall 2029",0.4,6.88,12.4,0.38,sz=10,bold=True,col=ACC)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — PROFESSIONAL DEVELOPMENT & GOALS
# ══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLK)
rect(s6,0,0,13.3,7.5,BG)
rect(s6,0,0,13.3,0.07,ACC)
txt(s6,"PROFESSIONAL DEVELOPMENT & IDP GOALS",0.3,0.1,9,0.35,sz=18,bold=True,col=WHT,font="Trebuchet MS")

goals = [
    ("GOAL 1: ACADEMIC PUBLICATIONS", PHASE1, [
        "Timeline: Phase 1-4 publications integrated into research phases",
        "Target: 5+ journal articles + 5+ conference papers by graduation",
        "Plan: Phase 1 (1 paper) → Phase 2 (1-2) → Phase 3 (2-3) → Phase 4 (1-2 final)",
        "Status: ✓ On track (publications embedded in Gantt)",
        ""
    ]),
    ("GOAL 2: TECHNICAL MASTERY", PHASE2, [
        "Deep Learning: CS 59300 (Vision), CS 57800 (ML), ECE 57000 (AI)",
        "Time Series: STAT 52000, STAT 52400, STAT 51400",
        "Optimization: CS 52000, CS 58000",
        "Target: 300 LeetCode problems + Cloud cert by Summer 2027",
        "Status: ✓ Courses front-loaded; co-op summer for coding sprint"
    ]),
    ("GOAL 3: INDUSTRY EXPERIENCE", PHASE3, [
        "Co-op Internship: Summer 2027 (3 months, Jun-Aug)",
        "Placement: Strategic gap between foundation & prelim phases",
        "Focus: Industry tech skills (cloud ML, DevOps, data engineering)",
        "Future: Internship opportunities in Year 3-4 as time permits",
        "Status: ✓ Scheduled; aligns with post-co-op CSE courses (Fa27 onward)"
    ]),
]

y_pos = 0.6
for (goal_title, color, goal_items) in goals:
    # Goal title
    rect(s6, 0.3, y_pos, 12.7, 0.25, color)
    txt(s6, goal_title, 0.45, y_pos+0.03, 12.4, 0.2, sz=10, bold=True, col=WHT)
    y_pos += 0.28
    
    # Goal details
    for item in goal_items:
        if item == "":
            y_pos += 0.08
        else:
            rect(s6, 0.3, y_pos, 12.7, 0.12, DK2, BDR, 0.2)
            txt(s6, item, 0.45, y_pos+0.01, 12.55, 0.11, sz=7, col=SLV, wrap=True)
            y_pos += 0.13

# Save
out = str(Path(__file__).with_name("GanttChart_Kamrul_PhD_ENHANCED_TEMP.pptx"))
prs.save(out)
# Rename to final name
import shutil, os, time
final_out = str(Path(__file__).with_name("GanttChart_Kamrul_PhD_ENHANCED.pptx"))
try:
    if Path(final_out).exists():
        for attempt in range(5):
            try:
                Path(final_out).unlink()
                break
            except:
                time.sleep(0.5)
except:
    pass
shutil.move(out, final_out)
print("PPTX regenerated successfully")
print(f"Saved: {out}")

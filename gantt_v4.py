from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def rect(slide, x, y, w, h, fill, line=None, lw=0.4):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape


def txt(
    slide,
    text,
    x,
    y,
    w,
    h,
    sz=10,
    bold=False,
    col="FFFFFF",
    align=PP_ALIGN.LEFT,
    font="Calibri",
    wrap=True,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = wrap
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = valign
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = rgb(col)
    return box


def diamond(slide, x, y, size, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def connector(slide, x1, y1, x2, y2, color, width=0.75, dash=True):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = 4
    return line


BG = "0B0F1A"
ACC = "3B82F6"
WHT = "FFFFFF"
SLV = "CBD5E1"
MUTED = "94A3B8"
DIM = "64748B"
BDR = "1E293B"
DK1 = "0D1220"
DK2 = "111827"
DK3 = "0F172A"
TRN = "4F46E5"
CUR = "2563EB"
COURSE = "7C3AED"
ADMIN = "0F766E"
PRELIM = "F59E0B"
RESEARCH = "BE185D"
PROPOSAL = "DC2626"
COMPLETE = "059669"
CONTINGENCY = "94A3B8"
SEMINAR = "1E3A5F"


prs = Presentation()
prs.slide_width = Inches(13.3)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, 13.3, 7.5, BG)
rect(s1, 0, 0, 13.3, 0.07, ACC)
rect(s1, 0, 0.07, 0.35, 7.43, DK1)

txt(
    s1,
    "MY GRADUATE PROGRAM TIMELINE",
    0.6,
    1.04,
    11.4,
    0.72,
    sz=31,
    bold=True,
    col=WHT,
    font="Trebuchet MS",
    wrap=False,
)
txt(
    s1,
    "Version 4 - Handbook-aligned prelim dependencies - Spring 2026 to Fall 2029",
    0.6,
    1.82,
    11.5,
    0.36,
    sz=13,
    col=SLV,
    wrap=False,
)
txt(
    s1,
    "Md Kamruzzaman Kamrul",
    0.6,
    2.45,
    7.5,
    0.34,
    sz=14,
    bold=True,
    col=ACC,
    wrap=False,
)
txt(
    s1,
    "PhD in Technology - Purdue Polytechnic Institute - CSE Concentration",
    0.6,
    2.83,
    9.5,
    0.3,
    sz=11.5,
    col=SLV,
    wrap=False,
)
txt(
    s1,
    "Advisor: Dr. Heung Jin Oh - Wii Lab - Bowen School of Construction Management Technology",
    0.6,
    3.16,
    10.4,
    0.3,
    sz=10.5,
    col=MUTED,
    wrap=False,
)

metrics = [
    ("103", "TOTAL CREDITS"),
    ("21", "TRANSFER (GSU)"),
    ("41", "PURDUE COURSEWORK"),
    ("41", "DISSERTATION"),
    ("4 YRS", "TARGET DURATION"),
]
for i, (value, label) in enumerate(metrics):
    x = 0.6 + i * 2.42
    rect(s1, x, 4.15, 2.2, 1.18, DK2, BDR, 0.5)
    txt(
        s1,
        value,
        x,
        4.26,
        2.2,
        0.52,
        sz=28,
        bold=True,
        col=ACC,
        align=PP_ALIGN.CENTER,
        font="Trebuchet MS",
        wrap=False,
    )
    txt(s1, label, x, 4.84, 2.2, 0.28, sz=8.6, col=MUTED, align=PP_ALIGN.CENTER)

legend = [
    (TRN, "Transfer"),
    (CUR, "Current"),
    (COURSE, "Coursework"),
    (ADMIN, "Admin/Prereq"),
    (PRELIM, "Prelim"),
    (RESEARCH, "Research"),
]
for i, (color, label) in enumerate(legend):
    x = 0.6 + i * 2.03
    rect(s1, x, 6.18, 0.22, 0.22, color)
    txt(s1, label, x + 0.28, 6.15, 1.65, 0.26, sz=9.2, col=SLV, wrap=False)


# ---------------------------------------------------------------------------
# Slide 2: Main Gantt chart
# ---------------------------------------------------------------------------
s2 = prs.slides.add_slide(BLANK)
rect(s2, 0, 0, 13.3, 7.5, BG)
rect(s2, 0, 0, 13.3, 0.07, ACC)
txt(
    s2,
    "VERSION 4 GANTT CHART: COURSEWORK, PRELIM DEPENDENCIES, RESEARCH",
    0.3,
    0.1,
    12.0,
    0.32,
    sz=15.5,
    bold=True,
    col=WHT,
    font="Trebuchet MS",
    wrap=False,
)
txt(
    s2,
    "Handbook checkpoints added: approved POS, committee, reading list, written/oral prelim, Form 10, retake window",
    0.3,
    0.48,
    12.4,
    0.23,
    sz=8.8,
    col=MUTED,
    wrap=False,
)

SEMS = [
    "Sp'26",
    "Su'26",
    "Fa'26",
    "Sp'27",
    "Su'27",
    "Fa'27",
    "Sp'28",
    "Su'28",
    "Fa'28",
    "Sp'29",
    "Su'29",
    "Fa'29",
]

# label, group, start semester index, span in semesters, color
ROWS = [
    ("GSU MS transfer credits - 7 courses / 21 cr", "TRANSFER", 0.00, 1.00, TRN),
    ("Sp'26 current: TECH603(1), STAT520(3), CS593(3), ECE595(3)", "COURSEWORK", 0.00, 1.00, CUR),
    ("Fa'26 core: MET527(3), STAT501(3), TECH601(1), TECH646(3)", "COURSEWORK", 2.00, 1.00, COURSE),
    ("Sp'27 methods/ML: TECH676(3), CS578(3), CS520(3)", "COURSEWORK", 3.00, 1.00, COURSE),
    ("Su'27 compressed electives: STAT514(3), CS580(3)", "COURSEWORK", 4.00, 1.00, COURSE),
    ("Fa'27 final prelim-semester courses: STAT524(3), ECE570(3)", "COURSEWORK", 5.00, 1.00, COURSE),
    ("GRAD689 CSE seminar - 0 cr recurring", "COURSEWORK", 5.00, 7.00, SEMINAR),
    ("Advisor topic/timeline meetings; GPA/core grade check", "PRELIM DEPENDENCIES", 3.00, 1.60, ADMIN),
    ("Plan of Study + prelim committee finalized (3+ faculty)", "PRELIM DEPENDENCIES", 3.30, 1.50, ADMIN),
    ("Develop reading list from research topic (10+ weeks out)", "PRELIM DEPENDENCIES", 4.00, 0.70, ADMIN),
    ("Submit reading list 9 weeks out; committee review/revise", "PRELIM DEPENDENCIES", 4.65, 0.55, PRELIM),
    ("Questions finalized + dept review; written component", "PRELIM DEPENDENCIES", 5.10, 0.45, PRELIM),
    ("Written review, oral exam, official result, Form 10", "PRELIM DEPENDENCIES", 5.55, 0.40, PRELIM),
    ("Retake window if needed - next immediate semester", "CONTINGENCY", 6.00, 1.00, CONTINGENCY),
    ("Research foundation: literature, problem framing, pilot ideas", "RESEARCH", 0.00, 4.70, RESEARCH),
    ("Prelim synthesis + proposal outline", "RESEARCH", 5.00, 1.00, RESEARCH),
    ("Proposal development: design, methods, committee feedback", "RESEARCH", 6.00, 2.70, PROPOSAL),
    ("Data collection and analysis", "RESEARCH", 9.00, 2.00, RESEARCH),
    ("Dissertation writing, revisions, defense preparation", "RESEARCH", 10.00, 2.00, COMPLETE),
]

LBL_W = 3.35
GROUP_W = 0.48
CHART_X = LBL_W
CHART_W = 13.3 - LBL_W - 0.12
CELL_W = CHART_W / len(SEMS)
HDR_Y = 0.80
HDR_H = 0.29
ROW_Y0 = HDR_Y + HDR_H + 0.04
ROW_H = 0.178
ROW_GAP = 0.011
GRP_GAP = 0.052

row_ys = []
cy = ROW_Y0
last_group = None
for _, group, *_ in ROWS:
    if last_group is not None and group != last_group:
        cy += GRP_GAP
    row_ys.append(cy)
    cy += ROW_H + ROW_GAP
    last_group = group
total_rows_h = cy - ROW_Y0

groups = {}
for idx, (_, group, *_rest) in enumerate(ROWS):
    groups.setdefault(group, []).append(idx)

for i, sem in enumerate(SEMS):
    x = CHART_X + i * CELL_W
    is_summer = sem.startswith("Su")
    is_current = i == 0
    fill = "090D18" if is_summer else ("0B0F1A" if i % 2 == 0 else "0D1322")
    rect(s2, x, ROW_Y0, CELL_W, total_rows_h + 0.03, fill)
    if is_current:
        rect(s2, x, ROW_Y0, CELL_W, total_rows_h + 0.03, "0B1E38")
    hfill = ACC if is_current else ("1C2D3D" if is_summer else BDR)
    rect(s2, x + 0.01, HDR_Y + 0.02, CELL_W - 0.02, HDR_H - 0.04, hfill)
    txt(
        s2,
        sem,
        x,
        HDR_Y + 0.035,
        CELL_W,
        HDR_H - 0.05,
        sz=7.3,
        bold=is_current,
        col=WHT if is_current else SLV,
        align=PP_ALIGN.CENTER,
        wrap=False,
    )

group_colors = {
    "TRANSFER": TRN,
    "COURSEWORK": COURSE,
    "PRELIM DEPENDENCIES": ADMIN,
    "CONTINGENCY": CONTINGENCY,
    "RESEARCH": RESEARCH,
}

drawn = set()
for idx, (label, group, start, span, color) in enumerate(ROWS):
    y = row_ys[idx]
    if group not in drawn:
        indices = groups[group]
        top = row_ys[indices[0]]
        bottom = row_ys[indices[-1]] + ROW_H
        h = bottom - top
        rect(s2, 0, top, GROUP_W, h, group_colors[group])
        rect(s2, GROUP_W, top, LBL_W - GROUP_W - 0.04, h, DK1)
        group_label = {
            "PRELIM DEPENDENCIES": "PRELIM\nDEPS",
            "CONTINGENCY": "CONTING",
        }.get(group, group)
        txt(
            s2,
            group_label,
            0.03,
            top + 0.03,
            GROUP_W - 0.06,
            h - 0.06,
            sz=5.1,
            bold=True,
            col=WHT,
            align=PP_ALIGN.CENTER,
            wrap=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        drawn.add(group)

    rect(s2, GROUP_W, y, LBL_W - GROUP_W - 0.04, ROW_H, DK1)
    txt(
        s2,
        label,
        GROUP_W + 0.05,
        y + 0.018,
        LBL_W - GROUP_W - 0.12,
        ROW_H - 0.03,
        sz=5.35,
        col=SLV,
        wrap=False,
    )

    bar_x = CHART_X + start * CELL_W + 0.014
    bar_w = span * CELL_W - 0.028
    rect(s2, bar_x, y + 0.017, bar_w, ROW_H - 0.034, color)

milestone_y = row_ys[-1] + ROW_H + 0.18
milestones = [
    (4.82, "POS +\ncourse audit", ADMIN, 0.00),
    (5.52, "Prelim\nExam", PRELIM, 0.14),
    (6.48, "Retake\nif needed", CONTINGENCY, 0.00),
    (8.50, "Proposal\nDefense", PROPOSAL, 0.14),
    (11.48, "Final\nDefense", PRELIM, 0.00),
]
for pos, label, color, label_offset in milestones:
    x = CHART_X + pos * CELL_W
    diamond(s2, x - 0.09, milestone_y, 0.18, color)
    connector(s2, x, ROW_Y0, x, milestone_y, color, width=0.7, dash=True)
    txt(
        s2,
        label,
        x - 0.43,
        milestone_y + 0.21 + label_offset,
        0.86,
        0.34,
        sz=6.2,
        bold=True,
        col=color,
        align=PP_ALIGN.CENTER,
        wrap=True,
    )

legend_y = 7.10
legend = [
    (TRN, "Transfer"),
    (CUR, "Current"),
    (COURSE, "Coursework"),
    (ADMIN, "Admin/Prereq"),
    (PRELIM, "Prelim"),
    (RESEARCH, "Research"),
    (CONTINGENCY, "Contingency"),
]
for i, (color, label) in enumerate(legend):
    x = 0.3 + i * 1.82
    rect(s2, x, legend_y + 0.04, 0.16, 0.15, color)
    txt(s2, label, x + 0.21, legend_y + 0.005, 1.42, 0.22, sz=7.7, col=SLV, wrap=False)


# ---------------------------------------------------------------------------
# Slide 3: Preliminary exam handbook timeline
# ---------------------------------------------------------------------------
s3 = prs.slides.add_slide(BLANK)
rect(s3, 0, 0, 13.3, 7.5, BG)
rect(s3, 0, 0, 13.3, 0.07, ACC)
txt(
    s3,
    "PRELIMINARY EXAM HANDBOOK TIMELINE",
    0.3,
    0.1,
    10.5,
    0.35,
    sz=18,
    bold=True,
    col=WHT,
    font="Trebuchet MS",
    wrap=False,
)
txt(
    s3,
    "Fall 2027 target with relative deadlines from the approved Polytechnic handbook",
    0.3,
    0.48,
    11.7,
    0.24,
    sz=9,
    col=MUTED,
    wrap=False,
)

cards = [
    (
        "1. Eligibility Gates",
        ADMIN,
        [
            "Approved Plan of Study.",
            "Good standing: 3.0+ GPA.",
            "No unresolved core grade below B-.",
            "Required courses complete or in progress.",
            "Committee: 3+ graduate faculty; majority regular.",
        ],
    ),
    (
        "2. Reading List",
        COURSE,
        [
            "Meet advisor on topic and timeline.",
            "Develop topic-linked list at least 10 weeks out.",
            "Submit to committee exactly 9 weeks out.",
            "Committee reviews/revises within 2 weeks.",
        ],
    ),
    (
        "3. Written Component",
        PRELIM,
        [
            "Questions finalized from approved list.",
            "Chair sends questions for department review.",
            "Complete written responses within 6 weeks.",
            "Submit responses immediately on completion.",
        ],
    ),
    (
        "4. Oral + Result",
        COMPLETE,
        [
            "Committee reviews written work within 2 weeks.",
            "Written result sent within 1 week.",
            "If approved, oral exam scheduled within 2 weeks.",
            "If both approved, Chair submits electronic Form 10.",
        ],
    ),
]

card_x = [0.35, 3.55, 6.75, 9.95]
card_w = 2.92
for (x, (title, color, items)) in zip(card_x, cards):
    rect(s3, x, 0.9, card_w, 0.34, color)
    txt(s3, title, x + 0.08, 0.955, card_w - 0.16, 0.22, sz=8.6, bold=True, col=WHT, wrap=False)
    y = 1.28
    for item in items:
        rect(s3, x, y, card_w, 0.36, DK2, BDR, 0.25)
        txt(s3, "- " + item, x + 0.08, y + 0.04, card_w - 0.16, 0.25, sz=6.85, col=SLV, wrap=True)
        y += 0.39

rect(s3, 0.35, 3.65, 12.6, 0.34, DK3, BDR, 0.4)
txt(s3, "Fall 2027 dependency chain for the Gantt milestone", 0.5, 3.71, 8.5, 0.22, sz=10, bold=True, col=ACC, wrap=False)

timeline = [
    ("-10 wks", "Reading list draft"),
    ("-9 wks", "Submit list"),
    ("-7 wks", "List approved"),
    ("-6 wks", "Questions to dept"),
    ("Exam", "Written component"),
    ("+2 wks", "Committee review"),
    ("+4 wks", "Oral exam"),
    ("+1 wk", "Form 10"),
]
start_x = 0.55
step_w = 1.52
line_y = 4.62
connector(s3, start_x + 0.2, line_y, start_x + step_w * 7 + 0.2, line_y, PRELIM, width=1.0, dash=False)
for i, (offset, label) in enumerate(timeline):
    x = start_x + i * step_w
    diamond(s3, x + 0.09, line_y - 0.1, 0.20, PRELIM if i != 7 else COMPLETE)
    txt(s3, offset, x - 0.22, line_y + 0.18, 0.82, 0.2, sz=7.1, bold=True, col=PRELIM, align=PP_ALIGN.CENTER, wrap=False)
    txt(s3, label, x - 0.34, line_y + 0.43, 1.10, 0.45, sz=6.25, col=SLV, align=PP_ALIGN.CENTER, wrap=True)

notes = [
    ("Pass Path", "Written and oral components must both be Approved in the same administration."),
    ("Retake Rule", "If Not Approved, the second prelim should be completed by Spring 2028; summer does not count as the next immediate semester."),
    ("Defense Guardrail", "Final defense in Fall 2029 leaves far more than two registration sessions after the Fall 2027 prelim."),
]
for i, (head, body) in enumerate(notes):
    x = 0.35 + i * 4.2
    rect(s3, x, 6.28, 4.0, 0.28, [COMPLETE, CONTINGENCY, ACC][i])
    txt(s3, head, x + 0.08, 6.325, 3.84, 0.2, sz=8, bold=True, col=WHT, wrap=False)
    rect(s3, x, 6.56, 4.0, 0.48, DK2, BDR, 0.25)
    txt(s3, body, x + 0.09, 6.61, 3.82, 0.38, sz=6.8, col=SLV, wrap=True)


# ---------------------------------------------------------------------------
# Slide 4: Milestone requirements and rationale
# ---------------------------------------------------------------------------
s4 = prs.slides.add_slide(BLANK)
rect(s4, 0, 0, 13.3, 7.5, BG)
rect(s4, 0, 0, 13.3, 0.07, ACC)
txt(
    s4,
    "KEY MILESTONES: REQUIREMENTS, DEPENDENCIES, OUTPUTS",
    0.3,
    0.1,
    11.5,
    0.35,
    sz=17,
    bold=True,
    col=WHT,
    font="Trebuchet MS",
    wrap=False,
)

milestone_rows = [
    (
        "AUG 2027 - POS + COURSEWORK AUDIT",
        ADMIN,
        "Approved POS, committee plan, GPA/core-grade review, and final prelim-semester course plan locked.",
        "Confirms handbook eligibility before the reading-list clock starts.",
    ),
    (
        "FALL 2027 - PRELIMINARY EXAM",
        PRELIM,
        "Reading list approved; written response completed; oral component passed; Chair submits electronic Form 10.",
        "Moves the plan from coursework-heavy preparation into proposal development.",
    ),
    (
        "SPRING 2028 - RETAKE WINDOW IF NEEDED",
        CONTINGENCY,
        "Second prelim only if written or oral component is Not Approved in Fall 2027.",
        "Handbook requires the second administration by the next immediate semester; summer does not apply.",
    ),
    (
        "FALL 2028 - PROPOSAL DEFENSE",
        PROPOSAL,
        "Finalize research design, methods, data plan, committee feedback, and IRB/data protocol if applicable.",
        "Creates a clean handoff into Spring-Summer 2029 data collection and analysis.",
    ),
    (
        "FALL 2029 - FINAL DEFENSE + ETD",
        COMPLETE,
        "Complete dissertation, defend, address revisions, submit ETD, and clear graduation requirements.",
        "Final milestone stays at the end and satisfies the post-prelim registration-session spacing rule.",
    ),
]

y = 0.65
for title, color, requirement, rationale in milestone_rows:
    rect(s4, 0.3, y, 12.7, 0.30, color)
    txt(s4, title, 0.45, y + 0.05, 12.4, 0.2, sz=9.3, bold=True, col=WHT, wrap=False)
    y += 0.32
    rect(s4, 0.3, y, 2.0, 0.35, DK3, BDR, 0.25)
    txt(s4, "Requirement", 0.4, y + 0.085, 1.8, 0.18, sz=7.2, bold=True, col=ACC, wrap=False)
    rect(s4, 2.35, y, 10.65, 0.35, DK2, BDR, 0.25)
    txt(s4, requirement, 2.45, y + 0.055, 10.4, 0.24, sz=7.2, col=SLV, wrap=True)
    y += 0.37
    rect(s4, 0.3, y, 2.0, 0.35, DK3, BDR, 0.25)
    txt(s4, "Why placed here", 0.4, y + 0.085, 1.8, 0.18, sz=7.2, bold=True, col=ACC, wrap=False)
    rect(s4, 2.35, y, 10.65, 0.35, DK2, BDR, 0.25)
    txt(s4, rationale, 2.45, y + 0.055, 10.4, 0.24, sz=7.2, col=SLV, wrap=True)
    y += 0.50


# ---------------------------------------------------------------------------
# Slide 5: Course inventory and credit breakdown
# ---------------------------------------------------------------------------
s5 = prs.slides.add_slide(BLANK)
rect(s5, 0, 0, 13.3, 7.5, BG)
rect(s5, 0, 0, 13.3, 0.07, ACC)
txt(
    s5,
    "COURSE INVENTORY & CREDIT BREAKDOWN",
    0.3,
    0.1,
    11.0,
    0.35,
    sz=18,
    bold=True,
    col=WHT,
    font="Trebuchet MS",
    wrap=False,
)
txt(
    s5,
    "Coursework is front-loaded so the Fall 2027 prelim has the handbook prerequisites in place.",
    0.3,
    0.48,
    12.0,
    0.24,
    sz=8.8,
    col=MUTED,
    wrap=False,
)

left_blocks = [
    (
        "TRANSFER - GSU MS (21 cr)",
        TRN,
        [
            "CIS 8005 Data Programming (3cr)",
            "CIS 8398 Advanced AI for Business (3cr)",
            "CIS 8695 Big Data Analytics (3cr)",
            "CIS 8795 Big Data Infrastructure (3cr)",
            "CIS 8080 IS Security & Privacy (3cr)",
            "CIS 8045 Unstructured Data Management (3cr)",
            "CIS 8690 Topics in Information Systems (3cr)",
        ],
    ),
    (
        "SPRING 2026 - Current (10 cr)",
        CUR,
        [
            "TECH 60300 Graduate Seminar - Planning (1cr)",
            "STAT 52000 Time Series & Applications (3cr)",
            "CS 59300 Computer Vision (3cr)",
            "ECE 59500 Selected Topics in ECE (3cr)",
        ],
    ),
    (
        "FALL 2026 - Core / Methods (10 cr)",
        COURSE,
        [
            "MET 52700 Technology: Global Perspective (3cr)",
            "STAT 50100 Experimental Statistics I (3cr)",
            "TECH 60100 Research Seminar in Technology (1cr)",
            "TECH 64600 Analysis of Research in Industry (3cr)",
        ],
    ),
]

right_blocks = [
    (
        "SPRING 2027 - Advanced Methods (9 cr)",
        COURSE,
        [
            "TECH 67600 Analysis of Research Methods (3cr)",
            "CS 57800 Statistical Machine Learning (3cr)",
            "CS 52000 Computational Optimization (3cr)",
        ],
    ),
    (
        "SUMMER 2027 - Final Prep Load (6 cr)",
        COURSE,
        [
            "STAT 51400 Design of Experiments (3cr)",
            "CS 58000 Algorithm Design & Analysis (3cr)",
            "Reading-list drafting and committee/POS audit",
        ],
    ),
    (
        "FALL 2027 - Final Coursework + Prelim (6 cr)",
        PRELIM,
        [
            "STAT 52400 Applied Multivariate Analysis (3cr)",
            "ECE 57000 Artificial Intelligence (3cr)",
            "GRAD 68900 CSE Seminar (0cr)",
            "Written prelim, oral prelim, Form 10",
        ],
    ),
    (
        "SPRING 2028 - FALL 2029 - Research Completion",
        RESEARCH,
        [
            "Sp'28: proposal writing; retake window only if needed",
            "Su'28-Fa'28: proposal completion and defense",
            "Sp'29-Su'29: data collection and analysis",
            "Fa'29: dissertation writing, final defense, ETD",
        ],
    ),
]


def draw_blocks(slide, blocks, x, y, w, item_h=0.184):
    for title, color, items in blocks:
        rect(slide, x, y, w, 0.24, color)
        txt(slide, title, x + 0.07, y + 0.035, w - 0.14, 0.17, sz=7.8, bold=True, col=WHT, wrap=False)
        y += 0.25
        for item in items:
            rect(slide, x, y, w, item_h, DK2, BDR, 0.25)
            txt(slide, item, x + 0.08, y + 0.025, w - 0.16, item_h - 0.04, sz=7.0, col=SLV, wrap=True)
            y += item_h + 0.006
        y += 0.06
    return y


draw_blocks(s5, left_blocks, 0.25, 0.78, 6.15, item_h=0.18)
draw_blocks(s5, right_blocks, 6.75, 0.78, 6.15, item_h=0.205)
rect(s5, 0.25, 6.82, 12.8, 0.52, DK3, BDR, 0.5)
txt(
    s5,
    "Credit summary: 21 transfer + 41 Purdue coursework + 41 dissertation/research credits = 103 credits listed | Target graduation: Fall 2029",
    0.4,
    6.91,
    12.4,
    0.27,
    sz=9.0,
    bold=True,
    col=ACC,
    align=PP_ALIGN.CENTER,
    wrap=False,
)


# ---------------------------------------------------------------------------
# Slide 6: Research phases, risk controls, and professional goals
# ---------------------------------------------------------------------------
s6 = prs.slides.add_slide(BLANK)
rect(s6, 0, 0, 13.3, 7.5, BG)
rect(s6, 0, 0, 13.3, 0.07, ACC)
txt(
    s6,
    "RESEARCH PHASES, CONTROLS, AND IDP ALIGNMENT",
    0.3,
    0.1,
    11.0,
    0.35,
    sz=17,
    bold=True,
    col=WHT,
    font="Trebuchet MS",
    wrap=False,
)

phase_cards = [
    (
        "PHASE 1: FOUNDATION",
        "Sp'26 - Sp'27",
        CUR,
        [
            "Literature review and research problem framing.",
            "Advisor meetings; first publication idea.",
            "Build technical base in ML, CV, statistics, optimization.",
        ],
    ),
    (
        "PHASE 2: QUALIFICATION",
        "Su'27 - Fa'27",
        PRELIM,
        [
            "POS/course audit, committee, reading list, written/oral prelim.",
            "Synthesis-focused study, not memorization-only review.",
            "Spring 2028 retake window exists only as contingency.",
        ],
    ),
    (
        "PHASE 3: PROPOSAL",
        "Sp'28 - Fa'28",
        PROPOSAL,
        [
            "Proposal writing, methodology, data protocol, committee feedback.",
            "IRB or data permissions if required by final design.",
            "Proposal defense in Fall 2028.",
        ],
    ),
    (
        "PHASE 4: COMPLETION",
        "Sp'29 - Fa'29",
        COMPLETE,
        [
            "Data collection and analysis in Spring-Summer 2029.",
            "Dissertation writing and revisions through Fall 2029.",
            "Final defense, ETD submission, graduation clearance.",
        ],
    ),
]

for i, (title, dates, color, items) in enumerate(phase_cards):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.55
    y = 0.72 + row * 2.56
    rect(s6, x, y, 6.1, 0.30, color)
    txt(s6, title, x + 0.1, y + 0.05, 4.3, 0.2, sz=8.7, bold=True, col=WHT, wrap=False)
    txt(s6, dates, x + 4.55, y + 0.05, 1.4, 0.2, sz=7.6, bold=True, col=WHT, align=PP_ALIGN.RIGHT, wrap=False)
    rect(s6, x, y + 0.32, 6.1, 1.50, DK2, BDR, 0.3)
    item_y = y + 0.42
    for item in items:
        txt(s6, "- " + item, x + 0.15, item_y, 5.8, 0.26, sz=7.2, col=SLV, wrap=True)
        item_y += 0.34

rect(s6, 0.35, 6.00, 12.6, 0.30, DK3, BDR, 0.4)
txt(s6, "Controls added in Version 4", 0.5, 6.055, 4.0, 0.2, sz=9.0, bold=True, col=ACC, wrap=False)
controls = [
    "Prelim is expanded into a checklist-driven workflow, not a single date.",
    "Retake contingency is visible but separated from the main success path.",
    "Research-only semesters are explicit: Sp'29, Su'29, and Fa'29.",
    "Final defense remains the terminal milestone with ETD/graduation clearance.",
]
for i, item in enumerate(controls):
    x = 0.45 + (i % 2) * 6.35
    y = 6.42 + (i // 2) * 0.36
    txt(s6, "- " + item, x, y, 6.0, 0.25, sz=7.2, col=SLV, wrap=True)


out = Path(__file__).with_name("GanttChart_Kamrul_PhD_v4.pptx")
prs.save(out)
print(f"Saved: {out}")
